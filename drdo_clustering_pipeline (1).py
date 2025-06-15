import os
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation
import pytesseract
from PIL import Image, ImageFilter
import re
from sentence_transformers import SentenceTransformer
import csv
from collections import defaultdict
import matplotlib.pyplot as plt
from sklearn.decomposition import PCA
from sklearn.cluster import KMeans
import numpy as np
import multiprocessing as mp
import shutil

INPUT_FOLDER = r"E:\CODING\Python\DRDO Internship Project\project_data"
OUTPUT_CSV = "output_embeddings_labels_multiprocessing.csv"
GROUPED_OUTPUT = "grouped_files_by_cluster_multiprocessing.txt"
OUTPUT_DIR = "DRDO_Files"

sbert_model = SentenceTransformer("all-MiniLM-L6-v2")

def extract_text_from_pdf(filepath):
    try:
        with fitz.open(filepath) as doc:
            return " ".join(page.get_text() for page in doc)
    except Exception as e:
        print(f"Error reading PDF {filepath}: {e}")
        return ""

def extract_text_from_docx(filepath):
    try:
        doc = docx.Document(filepath)
        return " ".join(para.text for para in doc.paragraphs)
    except Exception as e:
        print(f"Error reading DOCX {filepath}: {e}")
        return ""

def extract_text_from_excel(filepath):
    try:
        dfs = pd.read_excel(filepath, sheet_name=None, dtype=str)
        content = []
        for sheet_name, df in dfs.items():
            content.append(sheet_name)
            content.extend(df.columns.astype(str).tolist())
            for row in df.itertuples(index=False):
                content.append(" ".join(str(cell) for cell in row))
        return " ".join(content)
    except Exception as e:
        print(f"Error reading Excel {filepath}: {e}")
        return ""

def extract_text_from_pptx(filepath):
    try:
        prs = Presentation(filepath)
        return " ".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    except Exception as e:
        print(f"Error reading PPTX {filepath}: {e}")
        return ""

def extract_text_from_image(filepath):
    try:
        img = Image.open(filepath).convert("L").filter(ImageFilter.SHARPEN)
        return pytesseract.image_to_string(img)
    except Exception as e:
        print(f"Error reading Image {filepath}: {e}")
        return ""

def clean_text(text):
    text = text.lower()
    text = re.sub(r"[^a-z0-9\s]", "", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()

def extract_text(filepath):
    ext = filepath.lower().split(".")[-1]
    if ext == "pdf":
        return extract_text_from_pdf(filepath)
    elif ext == "docx":
        return extract_text_from_docx(filepath)
    elif ext in ["xlsx", "xls"]:
        return extract_text_from_excel(filepath)
    elif ext == "pptx":
        return extract_text_from_pptx(filepath)
    elif ext in ["jpg", "jpeg", "png"]:
        return extract_text_from_image(filepath)
    elif ext == "txt":
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            print(f"Error reading TXT {filepath}: {e}")
            return ""
    else:
        return ""

def process_file(filepath):
    text = extract_text(filepath)
    if text:
        cleaned = clean_text(text)
        return filepath, cleaned
    return None

if __name__ == "__main__":
    all_filepaths = []
    for root, dirs, files in os.walk(INPUT_FOLDER):
        for file in files:
            all_filepaths.append(os.path.join(root, file))

    print(f"✅ Found {len(all_filepaths)} files to process.")

    with mp.Pool(processes=4) as pool:  # Adjust number of processes based on CPU cores
        results = pool.map(process_file, all_filepaths)

    results = [res for res in results if res is not None]
    filepaths, texts = zip(*results)

    print(f"✅ Extracted and cleaned {len(texts)} documents.")

    print("🔄 Generating embeddings...")
    embeddings = sbert_model.encode(texts, show_progress_bar=True)

    print("📊 Visualizing embeddings with PCA...")
    pca = PCA(n_components=2)
    pca_embeddings = pca.fit_transform(embeddings)
    plt.figure(figsize=(8, 6))
    plt.scatter(pca_embeddings[:, 0], pca_embeddings[:, 1], alpha=0.7)
    plt.title("PCA Visualization of Embeddings")
    plt.xlabel("PC 1")
    plt.ylabel("PC 2")
    plt.show()

    print("🔄 Clustering embeddings with KMeans...")
    n_clusters = 4
    clusterer = KMeans(n_clusters=n_clusters, random_state=42)
    labels = clusterer.fit_predict(embeddings)

    plt.figure(figsize=(10, 8))
    palette = plt.cm.get_cmap('tab10', n_clusters)
    colors = [palette(l) for l in labels]
    plt.scatter(pca_embeddings[:, 0], pca_embeddings[:, 1], c=colors, s=40)
    plt.title(f"Clusters (KMeans, n_clusters={n_clusters})")
    plt.xlabel("PC 1")
    plt.ylabel("PC 2")
    plt.show()

    print(f"💾 Saving embeddings and cluster labels to: {OUTPUT_CSV}")
    with open(OUTPUT_CSV, "w", newline='', encoding="utf-8") as csvfile:
        writer = csv.writer(csvfile)
        writer.writerow(["filepath", "label"] + [f"dim_{i}" for i in range(embeddings.shape[1])])
        for filepath, label, vec in zip(filepaths, labels, embeddings):
            writer.writerow([filepath, label] + list(vec))

    clusters = defaultdict(list)
    for filepath, label in zip(filepaths, labels):
        clusters[label].append(filepath)

    print(f"💾 Saving grouped file list to: {GROUPED_OUTPUT}")
    with open(GROUPED_OUTPUT, "w", encoding="utf-8") as f:
        for label, files in sorted(clusters.items()):
            f.write(f"Cluster {label}:")
            f.write("\n")
            for file in files:
                f.write(f"  {file}\n")
            f.write("\n")

    print(f"📁 Creating folders and copying clustered files to '{OUTPUT_DIR}'...")
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    for label, files in sorted(clusters.items()):
        cluster_folder = os.path.join(OUTPUT_DIR, f"Cluster_{label}")
        os.makedirs(cluster_folder, exist_ok=True)
        for file in files:
            try:
                shutil.copy(file, cluster_folder)
            except Exception as e:
                print(f"Error copying file {file} to {cluster_folder}: {e}")

    print(f"✅ All clustered files copied into '{OUTPUT_DIR}' folder with subfolders by cluster.")
    print("✅ Process completed successfully!")
